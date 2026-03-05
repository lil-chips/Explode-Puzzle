#!/usr/bin/env python3
# -*- coding: utf-8 -*-

"""Generate CY Orchid (麒悅) wholesaler pitch deck in PPTX (EN + ZH-CN).

Notes:
- Uses local images from ./麒悅/*
- Produces 2 PPTX files with consistent layout.
"""

from __future__ import annotations

import os
from pathlib import Path
from typing import Iterable

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

ROOT = Path(__file__).resolve().parents[1]
ASSETS = ROOT / "麒悅"
OUTDIR = ASSETS / "Deck"
OUTDIR.mkdir(parents=True, exist_ok=True)


def rgb(hexstr: str) -> RGBColor:
    hexstr = hexstr.lstrip("#")
    return RGBColor(int(hexstr[0:2], 16), int(hexstr[2:4], 16), int(hexstr[4:6], 16))


BRAND_BROWN = rgb("4A2E1A")
BG_CREAM = rgb("F5EFE6")
TEXT_DARK = rgb("1B1B1B")
ACCENT = rgb("D83B72")


def pick_first(paths: Iterable[Path]) -> Path | None:
    for p in paths:
        if p.exists():
            return p
    return None


def add_title(slide, title: str, subtitle: str | None = None):
    left, top, width, height = Inches(0.7), Inches(0.5), Inches(12.0), Inches(1.0)
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.size = Pt(38)
    run.font.bold = True
    run.font.color.rgb = BRAND_BROWN

    if subtitle:
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.size = Pt(18)
        p2.font.color.rgb = rgb("555555")


def add_bullets(slide, heading: str, bullets: list[str], *, x=0.9, y=1.6, w=6.2, h=5.2):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()

    p0 = tf.paragraphs[0]
    p0.text = heading
    p0.font.size = Pt(24)
    p0.font.bold = True
    p0.font.color.rgb = BRAND_BROWN

    for b in bullets:
        p = tf.add_paragraph()
        p.text = b
        p.level = 0
        p.font.size = Pt(16)
        p.font.color.rgb = TEXT_DARK


def add_image(slide, img: Path, *, x, y, w=None, h=None):
    if not img or not img.exists():
        return
    if w is not None and h is not None:
        slide.shapes.add_picture(str(img), Inches(x), Inches(y), Inches(w), Inches(h))
    elif w is not None:
        slide.shapes.add_picture(str(img), Inches(x), Inches(y), width=Inches(w))
    elif h is not None:
        slide.shapes.add_picture(str(img), Inches(x), Inches(y), height=Inches(h))
    else:
        slide.shapes.add_picture(str(img), Inches(x), Inches(y))


def set_bg(slide, color: RGBColor):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_footer(slide, text: str):
    box = slide.shapes.add_textbox(Inches(0.7), Inches(6.95), Inches(12.0), Inches(0.4))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(11)
    p.font.color.rgb = rgb("666666")


def build_deck(lang: str):
    # lang: "en" or "zh"
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    logo = pick_first([
        ASSETS / "麒悅logo 2.png",
        ASSETS / "麒悅logo.jpg",
    ])

    company_imgs = sorted((ASSETS / "公司照片").glob("*.png")) + sorted((ASSETS / "公司照片").glob("*.JPG"))
    hero_factory = pick_first([ASSETS / "公司照片" / "溫室栽培區.png", ASSETS / "公司照片" / "包裝預冷區.png"])

    def products(folder: str, count: int = 4) -> list[Path]:
        d = ASSETS / folder
        # Prefer the plain code image (not p/f) if available
        imgs = sorted([p for p in d.glob("*.png") if " p." not in p.name and " f." not in p.name])
        if len(imgs) < count:
            imgs += sorted(d.glob("*.png"))
        # de-dup
        seen = set()
        out = []
        for p in imgs:
            if p.name in seen:
                continue
            seen.add(p.name)
            out.append(p)
            if len(out) >= count:
                break
        return out

    mini = products("迷你花種", 4)
    medium = products("中輪花種", 4)
    large = products("大輪花種", 4)

    # --- Slide 1: Cover
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, BG_CREAM)
    title = "CY Orchid — Taiwan Phalaenopsis Supply" if lang == "en" else "麒悦兰花（CY Orchid）— 台湾蝴蝶兰供货"
    subtitle = "For Melbourne / Sydney wholesalers" if lang == "en" else "面向墨尔本 / 悉尼批发商"
    add_title(slide, title, subtitle)
    if hero_factory:
        add_image(slide, hero_factory, x=7.4, y=1.2, w=5.3, h=5.5)
    if logo:
        add_image(slide, logo, x=0.7, y=6.2, w=2.0)
    add_footer(slide, "www.cyorchid.com.tw")

    # --- Slide 2: One-liner
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, BG_CREAM)
    h = "Why CY Orchid" if lang == "en" else "为什么选择麒悦"
    bullets = [
        "Established in 2002 • Pingtung, Taiwan" if lang == "en" else "2002年成立 • 台湾屏东",
        "Phalaenopsis specialist: flasks → greenhouse → packaging/cold-chain" if lang == "en" else "专注蝴蝶兰：瓶苗→温室→包装预冷/冷链",
        "Export experience: worldwide markets incl. Australia & NZ" if lang == "en" else "全球出口经验：含澳洲/纽西兰",
        "Sedex certified (CSR supply-chain standard)" if lang == "en" else "Sedex认证（供应链/劳工/ESG标准）",
    ]
    add_bullets(slide, h, bullets, x=0.9, y=0.9, w=7.0, h=5.8)
    add_image(slide, company_imgs[0] if company_imgs else None, x=8.2, y=1.0, w=4.8, h=5.6)
    add_footer(slide, "Taiwan-grown • Consistent quality • Wholesale-ready" if lang == "en" else "台湾产地 • 品质稳定 • 面向批发合作")

    # --- Slide 3: Product forms
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, BG_CREAM)
    h = "Product Forms" if lang == "en" else "供货形态"
    bullets = [
        "Flowering potted plants (ready for retail)" if lang == "en" else "开花盆栽（可直接零售）",
        "Young plants / flask seedlings (for growers)" if lang == "en" else "小苗/瓶苗（适合种植端）",
        "Custom order by market preferences (size/color)" if lang == "en" else "可按市场偏好客制（花型/色系/尺寸）",
    ]
    add_bullets(slide, h, bullets, x=0.9, y=0.9, w=7.0, h=5.8)
    add_image(slide, company_imgs[1] if len(company_imgs) > 1 else None, x=8.2, y=1.0, w=4.8, h=2.7)
    add_image(slide, company_imgs[2] if len(company_imgs) > 2 else None, x=8.2, y=3.95, w=4.8, h=2.7)
    add_footer(slide, "MOQ / lead time available upon request" if lang == "en" else "可提供MOQ与交期（按订单确认）")

    # --- Slides 4-6: Product categories
    def product_slide(title_txt: str, imgs: list[Path]):
        s = prs.slides.add_slide(prs.slide_layouts[6])
        set_bg(s, BG_CREAM)
        add_title(s, title_txt)
        # grid of 4 images
        positions = [(0.9, 1.6), (4.2, 1.6), (7.5, 1.6), (10.8, 1.6)]
        for (x, y), img in zip(positions, imgs):
            add_image(s, img, x=x, y=y, w=2.9, h=4.6)
        add_footer(s, "Variety codes available • Photos for reference" if lang == "en" else "提供品种编号 • 图片仅供参考")

    product_slide("Mini (迷你花)" if lang == "en" else "迷你花（Mini）", mini)
    product_slide("Medium (中轮花)" if lang == "en" else "中轮花（Medium）", medium)
    product_slide("Large (大轮花)" if lang == "en" else "大轮花（Large）", large)

    # --- Slide 7: Facilities & QC
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, BG_CREAM)
    h = "Facilities & Quality Control" if lang == "en" else "工厂与品控"
    bullets = [
        "Clean-room tissue culture lab + standardized greenhouse operations" if lang == "en" else "无尘组培实验室 + 标准化温室作业",
        "Packaging pre-cooling and cold-room flowering management" if lang == "en" else "包装预冷 + 冷房催花管理",
        "Quality checkpoints across propagation → cultivation → export prep" if lang == "en" else "从组培→栽培→出口全流程品控",
    ]
    add_bullets(slide, h, bullets, x=0.9, y=0.9, w=6.8, h=5.8)
    # collage
    add_image(slide, pick_first([ASSETS/"公司照片"/"組織培育操作區.png"]), x=7.9, y=0.95, w=5.1, h=2.05)
    add_image(slide, pick_first([ASSETS/"公司照片"/"包裝預冷區.png"]), x=7.9, y=3.2, w=2.5, h=2.5)
    add_image(slide, pick_first([ASSETS/"公司照片"/"專業品管把關.png"]), x=10.5, y=3.2, w=2.5, h=2.5)
    add_footer(slide, "Consistent quality for wholesale programs" if lang == "en" else "适合批发长期供货计划")

    # --- Slide 8: Cooperation
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, BG_CREAM)
    h = "Wholesale Cooperation" if lang == "en" else "批发合作方式"
    bullets = [
        "Target: Melbourne / Sydney wholesale partners" if lang == "en" else "目标合作：墨尔本/悉尼批发商",
        "We can provide: product list, MOQ, lead time, export packing options" if lang == "en" else "可提供：品项清单、MOQ、交期、出口包装方案",
        "Next step: share your preferred product form + price range + delivery terms" if lang == "en" else "下一步：请告知您偏好的供货形态/价位区间/交货条件",
    ]
    add_bullets(slide, h, bullets, x=0.9, y=0.9, w=7.0, h=5.8)
    add_image(slide, hero_factory, x=8.2, y=1.0, w=4.8, h=5.6)
    add_footer(slide, "Email us for quotation & samples" if lang == "en" else "欢迎来信获取报价与样品")

    # --- Slide 9: Contact
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, BG_CREAM)
    title = "Contact" if lang == "en" else "联系方式"
    add_title(slide, title)

    contact_lines = [
        "CY Orchid (麒悅企業)" if lang == "en" else "麒悦企业（CY Orchid）",
        "No. 1-68, Nanjin Rd., Wanluan Township, Pingtung County, Taiwan" if lang == "en" else "台湾屏东县万峦乡南进路1-68号",
        "Tel: +886-8-781-5561" if lang == "en" else "电话：+886-8-781-5561",
        "Website: www.cyorchid.com.tw",
    ]
    add_bullets(slide, "", contact_lines, x=0.9, y=1.6, w=7.5, h=4.8)
    if logo:
        add_image(slide, logo, x=9.2, y=2.0, w=3.6)
    add_footer(slide, "Thank you" if lang == "en" else "感谢")

    out = OUTDIR / ("CY_Orchid_Wholesaler_Deck_EN.pptx" if lang == "en" else "CY_Orchid_Wholesaler_Deck_ZH-CN.pptx")
    prs.save(out)
    return out


def main():
    en = build_deck("en")
    zh = build_deck("zh")
    print(en)
    print(zh)


if __name__ == "__main__":
    main()
