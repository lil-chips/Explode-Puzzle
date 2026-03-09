#!/usr/bin/env python3
# -*- coding: utf-8 -*-

"""Generate CY Orchid (麒悅) V3 EN deck (proposal-style, credibility-first).

Design goals (per user feedback):
- Add detailed factory/process/credibility section (3–5+ slides) with photos + readable content.
- Do NOT repeat contact email on every slide; keep it mainly on final Contact slide.
- Cover + page 2 should not feel empty: strong positioning + clear offering scope.
- Overview should NOT list variety codes; describe capability (flask → young → flowering; mini/medium/large).
- Keep readable typography: not too small, not too light.

Inputs:
- Photos: ./麒悅/公司照片 and ./麒悅/重要資訊
- Logos: ./麒悅/麒悅logo 2.png and ./麒悅/麒悅logo.jpg
- Variety spec images: ./麒悅/{迷你花種,中輪花種,大輪花種}

Outputs:
- ./麒悅/DeckV3/CY_Orchid_Wholesaler_Proposal_V3_EN.pptx

Note: PDF export is intentionally omitted here due to PowerPoint AppleEvent timeouts (-1712).
"""

from __future__ import annotations

from dataclasses import dataclass
from pathlib import Path
from typing import Optional

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
BASE = ROOT / "Projects" / "麒悅"
OUTDIR = BASE / "DeckV3"
OUTDIR.mkdir(parents=True, exist_ok=True)


def rgb(h: str) -> RGBColor:
    h = h.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


# Palette: warm minimal, readable.
BG = rgb("F7F3EE")
INK = rgb("1F1F1F")
MUTED = rgb("4A4A4A")  # darker than V2 for readability
GOLD = rgb("B79B5B")
ACCENT = rgb("1F4D3A")  # deep green


SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


@dataclass
class Variety:
    code: str
    img: Path
    english_name: str
    flower_d_cm: Optional[str]
    height_cm: Optional[str]
    pot_cm: Optional[str]


def parse_specs_from_image(img: Path) -> dict:
    # Keep the same manual mapping used in V2.
    mapping = {
        "NBM525": dict(english_name="Phal. Chi Yueh Chilling Wind", flower="6", height="35", pot="6 / 9 / 12"),
        "CPS251": dict(english_name="Phal. Chi Yueh Bodhisattva", flower="8", height="55", pot="9 / 12"),
        "CWS196": dict(english_name="Phal. Ching Ann Antenna", flower="7", height="35", pot="9 / 12"),
        "CBM52": dict(english_name="Phal. Tulcan", flower="7", height="35", pot="9 / 12"),
        "CPM35": dict(english_name="Phal. Ching Ann Antenna", flower="7", height="35", pot="9 / 12"),
        "CPM68": dict(english_name="Phal. Chi Yueh Purity", flower="7", height="35", pot="9 / 12"),
        "CPM67": dict(english_name="Phal. Mayfair", flower="7", height="35", pot="9 / 12"),
        "CWM89": dict(english_name="Phal. Lioulin R Lip", flower="10", height="60", pot="9 / 12"),
        "CWM49": dict(english_name="Phal. Chi Yueh Purity", flower="7", height="35", pot="9 / 12"),
        "CGL94": dict(english_name="Phal. Chi Yueh Elf", flower="11", height="65", pot="12"),
        "CPL45": dict(english_name="CPL45", flower=None, height=None, pot=None),
    }

    code = img.stem.split()[0].replace("_", "").replace("-", "").upper()
    if code not in mapping and code.endswith("P"):
        code = code[:-1]
    return mapping.get(code, dict(english_name=code, flower=None, height=None, pot=None))


def add_bg(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG


def add_top_rule(slide, left: str, right: str | None = None):
    # thin line
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.65), Inches(0.52), Inches(12.05), Inches(0.02))
    line.fill.solid()
    line.fill.fore_color.rgb = GOLD
    line.line.fill.background()

    # left label
    box = slide.shapes.add_textbox(Inches(0.65), Inches(0.2), Inches(8.5), Inches(0.35))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = left
    p.font.size = Pt(12)
    p.font.color.rgb = MUTED

    if right:
        box2 = slide.shapes.add_textbox(Inches(8.8), Inches(0.2), Inches(3.9), Inches(0.35))
        tf2 = box2.text_frame
        tf2.clear()
        p2 = tf2.paragraphs[0]
        p2.text = right
        p2.font.size = Pt(12)
        p2.font.color.rgb = MUTED
        p2.alignment = PP_ALIGN.RIGHT


def add_footer_website(slide, website: str = "www.cyorchid.com.tw"):
    box = slide.shapes.add_textbox(Inches(0.65), Inches(7.15), Inches(12.05), Inches(0.3))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = website
    p.font.size = Pt(10)
    p.font.color.rgb = MUTED


def add_title_block(slide, title: str, subtitle: str | None = None):
    box = slide.shapes.add_textbox(Inches(0.9), Inches(1.1), Inches(11.6), Inches(1.4))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = INK

    if subtitle:
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.size = Pt(18)
        p2.font.color.rgb = MUTED
        p2.space_before = Pt(8)


def add_card(slide, x, y, w, h, title: str, body_lines: list[str]):
    # card background
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    card.fill.solid()
    card.fill.fore_color.rgb = rgb("FFFFFF")
    card.line.color.rgb = rgb("E5DED3")
    card.line.width = Pt(1)

    # title
    tbox = slide.shapes.add_textbox(x + Inches(0.3), y + Inches(0.2), w - Inches(0.6), Inches(0.5))
    tf = tbox.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = ACCENT

    # body
    bbox = slide.shapes.add_textbox(x + Inches(0.3), y + Inches(0.7), w - Inches(0.6), h - Inches(0.9))
    tf2 = bbox.text_frame
    tf2.clear()
    for i, line in enumerate(body_lines):
        p2 = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
        p2.text = line
        p2.font.size = Pt(13)
        p2.font.color.rgb = INK
        p2.space_after = Pt(3)


def add_full_bleed_photo(slide, img: Path, opacity_overlay: bool = True):
    # full bleed image
    slide.shapes.add_picture(str(img), Inches(0), Inches(0), width=SLIDE_W, height=SLIDE_H)
    if opacity_overlay:
        overlay = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_W, SLIDE_H)
        overlay.fill.solid()
        overlay.fill.fore_color.rgb = rgb("000000")
        overlay.fill.transparency = 0.55
        overlay.line.fill.background()


def add_logo(slide, logo_path: Path, x, y, w=None, h=None):
    if not logo_path.exists():
        return
    slide.shapes.add_picture(str(logo_path), x, y, width=w, height=h)


def add_cover(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Background: greenhouse photo if available.
    hero = BASE / "公司照片" / "溫室栽培區.png"
    if hero.exists():
        add_full_bleed_photo(slide, hero, opacity_overlay=True)
    else:
        add_bg(slide)

    # Logo (use provided non-original logo image)
    logo = BASE / "麒悅logo 2.png"
    add_logo(slide, logo, Inches(0.9), Inches(0.7), w=Inches(2.2))

    # Title
    box = slide.shapes.add_textbox(Inches(0.9), Inches(2.2), Inches(11.8), Inches(2.4))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = "CY Orchid"
    p.font.size = Pt(56)
    p.font.bold = True
    p.font.color.rgb = rgb("FFFFFF")

    p2 = tf.add_paragraph()
    p2.text = "Taiwan Phalaenopsis Supplier — Wholesale Proposal"
    p2.font.size = Pt(22)
    p2.font.color.rgb = rgb("FFFFFF")

    # 3 key points
    add_card(
        slide,
        Inches(0.9),
        Inches(5.4),
        Inches(11.9),
        Inches(1.55),
        "Highlights",
        [
            "Export-ready programs (AU/NZ-ready cold-chain preparation)",
            "Multiple supply forms: flask seedlings → young plants → flowering potted plants",
            "Mini / Medium / Large sizes available with consistent quality checks",
        ],
    )


def add_offering_scope(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_top_rule(slide, "CY Orchid", "What We Offer")

    add_title_block(slide, "What we offer")

    # Left: big bullets
    box = slide.shapes.add_textbox(Inches(0.95), Inches(2.3), Inches(7.1), Inches(4.8))
    tf = box.text_frame
    tf.clear()

    bullets = [
        "Product forms: flask seedlings / young plants / flowering potted plants",
        "Flower sizes: mini / medium / large Phalaenopsis",
        "Wholesale assortment programs & stable supply planning",
        "Specs, MOQ & lead time: available upon request",
    ]
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = b
        p.font.size = Pt(18)
        p.font.color.rgb = INK
        p.level = 0
        p.space_after = Pt(8)

    # Right: photos (use 2 images)
    pics = [
        BASE / "公司照片" / "瓶苗區.png",
        BASE / "公司照片" / "溫室栽培區.png",
    ]
    x = Inches(8.4)
    y = Inches(2.35)
    w = Inches(4.6)
    h = Inches(2.1)
    for idx, pth in enumerate([p for p in pics if p.exists()][:2]):
        slide.shapes.add_picture(str(pth), x, y + Inches(2.35) * idx, width=w, height=h)

    add_footer_website(slide)


def add_about(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_top_rule(slide, "CY Orchid", "About")

    slide_title = slide.shapes.add_textbox(Inches(0.95), Inches(1.05), Inches(11.6), Inches(0.8))
    tf = slide_title.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = "About CY Orchid"
    p.font.size = Pt(34)
    p.font.bold = True
    p.font.color.rgb = INK

    # Main photo
    img = BASE / "重要資訊" / "關於麒悅.png"
    if img.exists():
        slide.shapes.add_picture(str(img), Inches(0.95), Inches(1.85), width=Inches(7.2))

    # Credibility cards
    add_card(
        slide,
        Inches(8.4),
        Inches(1.85),
        Inches(4.6),
        Inches(1.55),
        "Export experience",
        [
            "Serving overseas markets (incl. Australia & New Zealand)",
            "Programs designed for wholesale distribution",
        ],
    )
    add_card(
        slide,
        Inches(8.4),
        Inches(3.55),
        Inches(4.6),
        Inches(1.55),
        "Reliable operations",
        [
            "Stage-by-stage quality checkpoints",
            "Cold-chain preparation for export shipments",
        ],
    )
    add_card(
        slide,
        Inches(8.4),
        Inches(5.25),
        Inches(4.6),
        Inches(1.75),
        "Communication",
        [
            "Specs / availability / quotations provided upon request",
            "Long-term supply planning supported",
        ],
    )

    add_footer_website(slide)


def add_facilities(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_top_rule(slide, "CY Orchid", "Facilities")
    add_title_block(slide, "Facilities overview", "From tissue culture to greenhouse & packing")

    # 4-image grid (big, readable)
    images = [
        BASE / "公司照片" / "組織培育操作區.png",
        BASE / "公司照片" / "組織備料分注區.png",
        BASE / "公司照片" / "溫室栽培區.png",
        BASE / "公司照片" / "辦公室.png",
    ]
    imgs = [p for p in images if p.exists()]
    positions = [
        (Inches(0.95), Inches(2.55)),
        (Inches(7.05), Inches(2.55)),
        (Inches(0.95), Inches(5.0)),
        (Inches(7.05), Inches(5.0)),
    ]
    for (x, y), pth in zip(positions, imgs[:4]):
        slide.shapes.add_picture(str(pth), x, y, width=Inches(5.9), height=Inches(2.25))

    add_footer_website(slide)


def add_production_flow(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_top_rule(slide, "CY Orchid", "Production Flow")

    add_title_block(slide, "Production flow", "A clear process from lab → greenhouse → shipping")

    # Left: a vertical flow using cards
    flow_cards = [
        ("Tissue culture", ["Clean handling & controlled propagation"]),
        ("Nursery & greenhouse", ["Growth management & consistency"]),
        ("Quality checkpoints", ["Sorting & inspection across stages"]),
        ("Pre-cooling & packing", ["Shipment preparation for cold-chain"]),
    ]
    y = Inches(2.55)
    for title, body in flow_cards:
        add_card(slide, Inches(0.95), y, Inches(6.5), Inches(1.05), title, body)
        y += Inches(1.15)

    # Right: relevant photo
    img = BASE / "重要資訊" / "麒悅簡介.png"
    if img.exists():
        slide.shapes.add_picture(str(img), Inches(7.7), Inches(2.55), width=Inches(5.3))

    add_footer_website(slide)


def add_quality_control(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_top_rule(slide, "CY Orchid", "Quality Control")
    add_title_block(slide, "Quality control", "Consistency matters for wholesale programs")

    img = BASE / "公司照片" / "專業品管把關.png"
    if img.exists():
        slide.shapes.add_picture(str(img), Inches(0.95), Inches(2.55), width=Inches(6.9))

    add_card(
        slide,
        Inches(8.05),
        Inches(2.55),
        Inches(4.95),
        Inches(2.2),
        "Focus areas",
        [
            "Health & uniformity checks",
            "Size grading & presentation",
            "Stable output for repeat orders",
        ],
    )
    add_card(
        slide,
        Inches(8.05),
        Inches(4.9),
        Inches(4.95),
        Inches(2.2),
        "For buyers",
        [
            "Clear specs shared upon request",
            "Wholesale planning support",
        ],
    )

    add_footer_website(slide)


def add_cold_chain(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_top_rule(slide, "CY Orchid", "Cold-chain & Packing")
    add_title_block(slide, "Cold-chain & export packing", "Pre-cooling and shipment preparation")

    img1 = BASE / "公司照片" / "包裝預冷區.png"
    img2 = BASE / "公司照片" / "親本室及冷房.png"

    # Two large photos stacked
    x = Inches(0.95)
    if img1.exists():
        slide.shapes.add_picture(str(img1), x, Inches(2.55), width=Inches(6.9), height=Inches(2.25))
    if img2.exists():
        slide.shapes.add_picture(str(img2), x, Inches(4.95), width=Inches(6.9), height=Inches(2.25))

    add_card(
        slide,
        Inches(8.05),
        Inches(2.55),
        Inches(4.95),
        Inches(4.65),
        "Designed for export",
        [
            "Pre-cooling prior to shipment",
            "Cold-room management & handling",
            "Packing prepared for downstream distribution",
            "Details provided upon request",
        ],
    )

    add_footer_website(slide)


def add_variety_slide(prs: Presentation, v: Variety):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_top_rule(slide, "CY Orchid", v.code)

    # Large photo
    slide.shapes.add_picture(str(v.img), Inches(0.95), Inches(1.1), width=Inches(7.25))

    # Right spec panel
    box = slide.shapes.add_textbox(Inches(8.35), Inches(1.1), Inches(4.65), Inches(5.8))
    tf = box.text_frame
    tf.clear()

    p0 = tf.paragraphs[0]
    p0.text = v.english_name
    p0.font.size = Pt(26)
    p0.font.bold = True
    p0.font.color.rgb = INK

    # divider
    div = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(8.35), Inches(1.78), Inches(4.65), Inches(0.02))
    div.fill.solid(); div.fill.fore_color.rgb = GOLD
    div.line.fill.background()

    def add_row(label: str, value: str):
        p = tf.add_paragraph()
        p.text = f"{label}: {value}"
        p.font.size = Pt(15)
        p.font.color.rgb = INK
        p.space_before = Pt(10)

    if v.flower_d_cm:
        add_row("Flower diameter", f"{v.flower_d_cm} cm")
        add_row("Plant height", f"{v.height_cm} cm")
        add_row("Recommended pot", f"{v.pot_cm} cm")
    else:
        p = tf.add_paragraph()
        p.text = "Specs available upon request"
        p.font.size = Pt(15)
        p.font.color.rgb = MUTED
        p.space_before = Pt(14)

    # Note: no email on this slide
    add_footer_website(slide)


def add_contact(prs: Presentation, email: str):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_top_rule(slide, "CY Orchid", "Contact")
    add_title_block(slide, "Contact")

    box = slide.shapes.add_textbox(Inches(0.95), Inches(2.4), Inches(11.7), Inches(3.2))
    tf = box.text_frame
    tf.clear()

    lines = [
        "CY Orchid (麒悅)",
        "Pingtung, Taiwan",
        f"Email: {email}",
        "Website: www.cyorchid.com.tw",
        "Specs / availability / quotations available upon request",
    ]

    p = tf.paragraphs[0]
    p.text = lines[0]
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = INK

    for t in lines[1:]:
        p2 = tf.add_paragraph()
        p2.text = t
        p2.font.size = Pt(18)
        p2.font.color.rgb = MUTED
        p2.space_before = Pt(6)


def build_en() -> Path:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    add_cover(prs)
    add_offering_scope(prs)

    # Credibility / factory section (3–5+ slides)
    add_about(prs)
    add_facilities(prs)
    add_production_flow(prs)
    add_quality_control(prs)
    add_cold_chain(prs)

    # Varieties (keep as in V2 for now)
    targets = [
        ("NBM525", BASE / "迷你花種" / "NBM525.png"),
        ("CPS251", BASE / "迷你花種" / "CPS251.png"),
        ("CWS196", BASE / "迷你花種" / "CWS196.png"),
        ("CBM52", BASE / "中輪花種" / "CBM52.png"),
        ("CPM35", BASE / "中輪花種" / "CPM35.png"),
        ("CPM68", BASE / "中輪花種" / "CPM68.png"),
        ("CPM67", BASE / "中輪花種" / "CPM67.png"),
        ("CWM89", BASE / "中輪花種" / "CWM89.png"),
        ("CWM49", BASE / "中輪花種" / "CWM49.png"),
        ("CGL94", BASE / "大輪花種" / "CGL94.png"),
        ("CPL45", BASE / "大輪花種" / "CPL45 p.png"),
    ]

    varieties: list[Variety] = []
    for code, path in targets:
        if not path.exists():
            continue
        spec = parse_specs_from_image(path)
        varieties.append(
            Variety(
                code=code,
                img=path,
                english_name=spec["english_name"],
                flower_d_cm=spec.get("flower"),
                height_cm=spec.get("height"),
                pot_cm=spec.get("pot"),
            )
        )

    for v in varieties:
        add_variety_slide(prs, v)

    add_contact(prs, email="skw10615@gmail.com")

    out = OUTDIR / "CY_Orchid_Wholesaler_Proposal_V3_EN.pptx"
    prs.save(out)
    return out


def main():
    pptx = build_en()
    print(pptx)


if __name__ == "__main__":
    main()
