#!/usr/bin/env python3
# -*- coding: utf-8 -*-

"""Generate CY Orchid (麒悅) V2 pitch deck in PPTX + PDF (EN + ZH-CN).

Design goals:
- Luxury minimal look: large photos, lots of whitespace, thin lines, restrained palette.
- One variety per slide (mostly), with specs extracted from provided images.

Inputs:
- Variety images in ./麒悅/{迷你花種,中輪花種,大輪花種}
- Company/process images in ./麒悅/重要資訊 and ./麒悅/公司照片

Outputs:
- ./麒悅/DeckV2/CY_Orchid_Catalog_V2_EN.pptx
- ./麒悅/DeckV2/CY_Orchid_Catalog_V2_ZH-CN.pptx
- plus PDFs via PowerPoint AppleScript
"""

from __future__ import annotations

import re
import subprocess
from dataclasses import dataclass
from pathlib import Path
from typing import Optional

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

ROOT = Path(__file__).resolve().parents[1]
BASE = ROOT / "Projects" / "麒悅"
OUTDIR = BASE / "DeckV2"
OUTDIR.mkdir(parents=True, exist_ok=True)

# --- palette

def rgb(h: str) -> RGBColor:
    h = h.lstrip('#')
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))

BG = rgb('F7F3EE')
INK = rgb('1F1F1F')
MUTED = rgb('6B6B6B')
GOLD = rgb('B79B5B')


@dataclass
class Variety:
    code: str
    img: Path
    english_name: str
    flower_d_cm: Optional[str]
    height_cm: Optional[str]
    pot_cm: Optional[str]


def ocr_text(img: Path) -> str:
    # macOS has built-in OCR via Vision? We'll use `tesseract` if available, else fallback to empty.
    # In this environment, we rely on the fact that model can read the images, so we parse from filename-based mapping.
    return ""


def parse_specs_from_image(img: Path) -> dict:
    """Extract specs from the image by simple regex on embedded text using OCR-like fallback.

    Since OCR availability varies, we use a heuristic: for these catalog images, the text is readable.
    We'll call `sips -g`? Not helpful. We'll instead maintain a manual mapping for required varieties.
    """
    # Manual mapping extracted from visual inspection (read tool) and consistent pattern.
    # Update if you add new varieties.
    mapping = {
        'NBM525': dict(english_name='Phal. Chi Yueh Chilling Wind', flower='6', height='35', pot='6 / 9 / 12'),
        'CPS251': dict(english_name='Phal. Chi Yueh Bodhisattva', flower='8', height='55', pot='9 / 12'),
        'CWS196': dict(english_name='Phal. Ching Ann Antenna', flower='7', height='35', pot='9 / 12'),
        'CBM52': dict(english_name='Phal. Tulcan', flower='7', height='35', pot='9 / 12'),
        'CPM35': dict(english_name='Phal. Ching Ann Antenna', flower='7', height='35', pot='9 / 12'),
        'CPM68': dict(english_name='Phal. Chi Yueh Purity', flower='7', height='35', pot='9 / 12'),
        'CPM67': dict(english_name='Phal. Mayfair', flower='7', height='35', pot='9 / 12'),
        'CWM89': dict(english_name='Phal. Lioulin R Lip', flower='10', height='60', pot='9 / 12'),
        'CWM49': dict(english_name='Phal. Chi Yueh Purity', flower='7', height='35', pot='9 / 12'),
        'CGL94': dict(english_name='Phal. Chi Yueh Elf', flower='11', height='65', pot='12'),
        'CPL45': dict(english_name='CPL45', flower=None, height=None, pot=None),
    }

    code = img.stem.split()[0].replace('_', '').replace('-', '')
    code = code.upper()
    if code not in mapping and code.endswith('P'):
        code = code[:-1]
    return mapping.get(code, dict(english_name=code, flower=None, height=None, pot=None))


def add_bg(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG


def add_brand_bar(slide, text_left: str, text_right: str):
    # Thin top line
    shape = slide.shapes.add_shape(1, Inches(0.6), Inches(0.45), Inches(12.1), Inches(0.02))
    shape.fill.solid()
    shape.fill.fore_color.rgb = GOLD
    shape.line.fill.background()

    # Left label
    box = slide.shapes.add_textbox(Inches(0.6), Inches(0.2), Inches(7), Inches(0.4))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text_left
    p.font.size = Pt(12)
    p.font.color.rgb = MUTED

    # Right label
    box2 = slide.shapes.add_textbox(Inches(8.0), Inches(0.2), Inches(4.7), Inches(0.4))
    tf2 = box2.text_frame
    tf2.clear()
    p2 = tf2.paragraphs[0]
    p2.text = text_right
    p2.font.size = Pt(12)
    p2.font.color.rgb = MUTED
    p2.alignment = PP_ALIGN.RIGHT


def add_title(slide, title: str, subtitle: str | None = None):
    box = slide.shapes.add_textbox(Inches(0.8), Inches(1.0), Inches(11.8), Inches(1.0))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = INK
    if subtitle:
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.size = Pt(18)
        p2.font.color.rgb = MUTED


def add_footer(slide, text: str):
    box = slide.shapes.add_textbox(Inches(0.6), Inches(7.1), Inches(12.1), Inches(0.3))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(10)
    p.font.color.rgb = MUTED


def add_variety_slide(prs: Presentation, v: Variety, lang: str, email: str):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_brand_bar(slide, "CY Orchid" if lang == 'en' else "麒悦兰花 CY Orchid", v.code)

    # Image big left
    slide.shapes.add_picture(str(v.img), Inches(0.8), Inches(1.1), width=Inches(7.2))

    # Right spec card
    box = slide.shapes.add_textbox(Inches(8.35), Inches(1.15), Inches(4.2), Inches(5.4))
    tf = box.text_frame
    tf.clear()

    name_title = v.english_name if lang == 'en' else v.code
    p0 = tf.paragraphs[0]
    p0.text = name_title
    p0.font.size = Pt(24)
    p0.font.bold = True
    p0.font.color.rgb = INK

    # thin divider
    div = slide.shapes.add_shape(1, Inches(8.35), Inches(1.75), Inches(4.2), Inches(0.02))
    div.fill.solid(); div.fill.fore_color.rgb = GOLD
    div.line.fill.background()

    def add_row(label_en, label_zh, value):
        p = tf.add_paragraph()
        label = label_en if lang == 'en' else label_zh
        p.text = f"{label}: {value}"
        p.font.size = Pt(14)
        p.font.color.rgb = INK
        p.space_before = Pt(10)

    if v.flower_d_cm:
        add_row("Flower diameter", "花径", f"{v.flower_d_cm} cm")
        add_row("Plant height", "株高", f"{v.height_cm} cm")
        add_row("Recommended pot", "适合盆径", f"{v.pot_cm} cm")
    else:
        p = tf.add_paragraph();
        p.text = "Specs available upon request" if lang == 'en' else "规格可另行提供"
        p.font.size = Pt(14)
        p.font.color.rgb = MUTED
        p.space_before = Pt(14)

    # CTA
    p = tf.add_paragraph()
    p.text = f"Contact: {email}"
    p.font.size = Pt(12)
    p.font.color.rgb = MUTED
    p.space_before = Pt(18)

    add_footer(slide, "www.cyorchid.com.tw")


def build(lang: str, email: str) -> Path:
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    # cover
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    title = "Premium Taiwan Phalaenopsis Catalog" if lang == 'en' else "台湾蝴蝶兰精品型录"
    subtitle = "For Melbourne / Sydney wholesalers" if lang == 'en' else "面向墨尔本 / 悉尼批发商"
    add_title(slide, title, subtitle)
    add_brand_bar(slide, "CY Orchid" if lang == 'en' else "麒悦兰花 CY Orchid", email)
    add_footer(slide, "www.cyorchid.com.tw")

    # overview
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_brand_bar(slide, "CY Orchid" if lang == 'en' else "麒悦兰花 CY Orchid", "Product Overview" if lang == 'en' else "产品总览")
    box = slide.shapes.add_textbox(Inches(0.8), Inches(1.0), Inches(11.8), Inches(5.6))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = "Product Lines" if lang == 'en' else "产品线"
    p.font.size = Pt(30)
    p.font.bold = True
    p.font.color.rgb = INK

    lines = [
        ("Mini", "迷你花", ["NBM525", "CPS251", "CWS196"]),
        ("Medium", "中轮花", ["CBM52", "CPM35", "CPM68", "CPM67", "CWM89", "CWM49"]),
        ("Large", "大轮花", ["CGL94", "CPL45"]),
    ]

    for name_en, name_zh, codes in lines:
        p2 = tf.add_paragraph()
        p2.text = (name_en if lang == 'en' else name_zh) + ": " + ", ".join(codes)
        p2.font.size = Pt(18)
        p2.font.color.rgb = INK
        p2.space_before = Pt(12)

    p3 = tf.add_paragraph()
    p3.text = ("We offer flowering potted plants, young plants and flask seedlings. Specs shown are for reference; availability may vary." if lang=='en'
               else "可供：开花盆栽 / 小苗 / 瓶苗。规格仅供参考；实际供货以订单确认。")
    p3.font.size = Pt(13)
    p3.font.color.rgb = MUTED
    p3.space_before = Pt(18)

    add_footer(slide, f"Contact: {email}")

    # varieties
    targets = [
        ("NBM525", BASE/"迷你花種"/"NBM525.png"),
        ("CPS251", BASE/"迷你花種"/"CPS251.png"),
        ("CWS196", BASE/"迷你花種"/"CWS196.png"),
        ("CBM52",  BASE/"中輪花種"/"CBM52.png"),
        ("CPM35",  BASE/"中輪花種"/"CPM35.png"),
        ("CPM68",  BASE/"中輪花種"/"CPM68.png"),
        ("CPM67",  BASE/"中輪花種"/"CPM67.png"),
        ("CWM89",  BASE/"中輪花種"/"CWM89.png"),
        ("CWM49",  BASE/"中輪花種"/"CWM49.png"),
        ("CGL94",  BASE/"大輪花種"/"CGL94.png"),
        ("CPL45",  BASE/"大輪花種"/"CPL45 p.png"),
    ]

    varieties: list[Variety] = []
    for code, path in targets:
        if not path.exists():
            continue
        spec = parse_specs_from_image(path)
        varieties.append(Variety(code=code, img=path, english_name=spec['english_name'],
                                flower_d_cm=spec.get('flower'), height_cm=spec.get('height'), pot_cm=spec.get('pot')))

    for v in varieties:
        add_variety_slide(prs, v, lang, email)

    # process / company
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_brand_bar(slide, "CY Orchid" if lang == 'en' else "麒悦兰花 CY Orchid", "Production & QC" if lang=='en' else "生产与品控")
    box = slide.shapes.add_textbox(Inches(0.8), Inches(1.0), Inches(6.4), Inches(5.8))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = "Production Flow" if lang=='en' else "生产流程"
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = INK

    bullets = [
        ("Tissue culture (clean-room) → greenhouse cultivation" if lang=='en' else "无尘组培 → 温室栽培"),
        ("Quality checkpoints across stages" if lang=='en' else "各阶段品质把关"),
        ("Packaging pre-cooling → cold-room management" if lang=='en' else "包装预冷 → 冷房管理"),
        ("Export-ready packing & cold-chain preparation" if lang=='en' else "出口包装与冷链准备"),
    ]
    for b in bullets:
        p2 = tf.add_paragraph(); p2.text = b
        p2.font.size = Pt(16); p2.font.color.rgb = INK
        p2.level = 0

    p3 = tf.add_paragraph();
    p3.text = ("CSR: Sedex certified" if lang=='en' else "CSR：Sedex认证")
    p3.font.size = Pt(13); p3.font.color.rgb = MUTED
    p3.space_before = Pt(12)

    # right image
    proc_img = BASE/"重要資訊"/"麒悅簡介.png"
    if proc_img.exists():
        slide.shapes.add_picture(str(proc_img), Inches(7.6), Inches(1.2), width=Inches(5.0))

    add_footer(slide, f"Contact: {email}  •  www.cyorchid.com.tw")

    # contact
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_brand_bar(slide, "CY Orchid" if lang=='en' else "麒悦兰花 CY Orchid", "Contact" if lang=='en' else "联系")
    add_title(slide, "Contact" if lang=='en' else "联系方式")
    box = slide.shapes.add_textbox(Inches(0.9), Inches(2.3), Inches(11.6), Inches(3.2))
    tf = box.text_frame
    tf.clear()
    lines = [
        ("CY Orchid (麒悅企業)" if lang=='en' else "麒悦企业（CY Orchid）"),
        ("Pingtung, Taiwan" if lang=='en' else "台湾屏东"),
        (f"Email: {email}"),
        ("Website: www.cyorchid.com.tw"),
    ]
    p = tf.paragraphs[0]
    p.text = lines[0]
    p.font.size = Pt(20); p.font.bold = True; p.font.color.rgb = INK
    for t in lines[1:]:
        p2 = tf.add_paragraph(); p2.text = t
        p2.font.size = Pt(16); p2.font.color.rgb = MUTED
        p2.space_before = Pt(6)

    out = OUTDIR / ("CY_Orchid_Catalog_V2_EN.pptx" if lang=='en' else "CY_Orchid_Catalog_V2_ZH-CN.pptx")
    prs.save(out)
    return out


def export_pdf(pptx: Path) -> Path:
    pdf = pptx.with_suffix('.pdf')
    script = ROOT / 'tools' / 'export_pptx_to_pdf.scpt'
    subprocess.run(['osascript', str(script), str(pptx), str(pdf)], check=True)
    return pdf


def main():
    email = 'skw10615@gmail.com'
    en = build('en', email)
    zh = build('zh', email)
    export_pdf(en)
    export_pdf(zh)
    print(en)
    print(zh)


if __name__ == '__main__':
    main()
